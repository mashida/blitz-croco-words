"""
Задача 2024.10.23.06

Разархивируйте архив и достаньте один файл презентации.

Откройте файл презентации в программе с помощью объекта Presentation библиотеки pptx
"""
from pathlib import Path
from typing import IO, Sequence

from pyaspeller import YandexSpeller
from pptx import Presentation

FILE_NAME = "Osennyaya_igra_3.pptx"


def is_not_valid(text: str) -> bool:
    return ' ' in text or '-' in text or ':' in text or 'СУПЕРКРОКО' in text


def get_words_form_file(file: IO[bytes]) -> list[str]:
    file_name = getattr(file, "name", "<stream>")
    print(f"getting words from {file_name} file")

    prs = Presentation(file)

    result: list[str] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if is_not_valid(shape.text):
                continue
            result.append(shape.text.strip())

    print(f"getting words done. got {len(result)} words")

    return result


def check_spelling(words: Sequence[str]) -> list[str]:
    if not words:
        return []

    print(f"checking {len(words)} words for spelling with Yandex Speller")

    speller = YandexSpeller()
    result = speller.spelled(' '.join(words))

    print(f"spell checking done")
    return result.split()


def save_words_to_file(words: list[str], filename: str | Path) -> None:
    path = Path(filename)
    print(f"saving {len(words)} words to {path} file")

    with path.open(mode='w', encoding='utf-8') as file:
        file.writelines([word + "\n" for word in words])

    print(f"saving words done")


if __name__ == '__main__':
    sample_path = Path(__file__).resolve().parent / FILE_NAME
    with sample_path.open('rb') as f:
        print(get_words_form_file(f))
